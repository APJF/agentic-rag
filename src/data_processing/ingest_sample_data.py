import os
import json
import re
from typing import Dict, List, Optional, Tuple
import argparse

from src.core.embedding import encode_text
from src.core.vector_store_interface import get_db_connection


def _read_docx(path: str) -> List[str]:
    try:
        from docx import Document  # python-docx
    except Exception as exc:
        raise RuntimeError("Thiếu dependency python-docx. Hãy cài đặt trước.") from exc

    doc = Document(path)
    paragraphs: List[str] = []
    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if text:
            paragraphs.append(text)
    return paragraphs


def _read_pptx(path: str) -> List[Tuple[int, str]]:
    try:
        from pptx import Presentation  # python-pptx
    except Exception as exc:
        raise RuntimeError("Thiếu dependency python-pptx. Hãy cài đặt trước.") from exc

    prs = Presentation(path)
    slides_text: List[Tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = (p.text or "").strip()
                    if line:
                        texts.append(line)
        combined = "\n".join(texts).strip()
        if combined:
            slides_text.append((idx, combined))
    return slides_text


def _chunk_text(lines: List[str], max_chars: int = 1200, overlap: int = 150) -> List[str]:
    if not lines:
        return []
    chunks: List[str] = []
    buf: List[str] = []
    current_len = 0
    for line in lines:
        if current_len + len(line) + 1 > max_chars and buf:
            chunks.append("\n".join(buf))
            # overlap
            if overlap > 0:
                carry = " \n".join(buf)[-overlap:]
                buf = [carry]
                current_len = len(carry)
            else:
                buf = []
                current_len = 0
        buf.append(line)
        current_len += len(line) + 1
    if buf:
        chunks.append("\n".join(buf))
    return chunks


def _normalize_text(text: str) -> str:
    # Chuẩn hóa đơn giản để giảm trùng lặp do khoảng trắng/line break
    text = (text or "").strip()
    # Gộp nhiều khoảng trắng thành một
    text = re.sub(r"\s+", " ", text)
    return text


def _infer_level_from_path(path: str) -> Optional[str]:
    path_upper = path.upper()
    for lv in ["N5", "N4", "N3", "N2", "N1"]:
        if lv in path_upper:
            return lv
    return None


def _infer_skill_from_path(path: str) -> Optional[str]:
    lower = path.lower()
    if any(k in lower for k in ["kanji"]):
        return "KANJI"
    if any(k in lower for k in ["vocab", "từ vựng", "tu vung"]):
        return "VOCAB"
    if any(k in lower for k in ["ngữ pháp", "ngu phap", "grammar"]):
        return "GRAMMAR"
    if any(k in lower for k in ["đọc", "doc", "reading"]):
        return "READING"
    if any(k in lower for k in ["nghe", "listening"]):
        return "LISTENING"
    if any(k in lower for k in ["nói", "noi", "speaking"]):
        return "SPEAKING"
    if any(k in lower for k in ["viết", "viet", "writing"]):
        return "WRITING"
    return None


def _insert_chunk(
    cur,
    table: str,
    chunk_text: str,
    source_name: str,
    original_page_number: Optional[int],
    level: Optional[str],
    skill: Optional[str],
    metadata: Dict
):
    # Dedup: kiểm tra tồn tại theo (source_document_name, original_page_number, chunk_text chuẩn hóa)
    n_text = _normalize_text(chunk_text)
    cur.execute(
        f"SELECT 1 FROM {table} WHERE source_document_name=%s AND original_page_number IS NOT DISTINCT FROM %s AND chunk_text=%s LIMIT 1;",
        (source_name, original_page_number, n_text),
    )
    if cur.fetchone():
        return  # đã có -> bỏ qua

    vector = encode_text(n_text)
    vector_str = str(list(vector))

    # Lấy danh sách cột có trong bảng để build INSERT động, tránh lỗi schema khác nhau
    cur.execute(
        "SELECT column_name FROM information_schema.columns WHERE table_name = %s;",
        (table,),
    )
    available_cols = {r[0] for r in cur.fetchall()}

    # Quy ước cột metadata
    metadata_col = None
    for candidate in ("metadata_json", "metadata", "meta_data"):
        if candidate in available_cols:
            metadata_col = candidate
            break

    columns = ["chunk_text", "source_document_name", "original_page_number"]
    values = [n_text, source_name, original_page_number]
    if "level" in available_cols:
        columns.append("level")
        values.append(level)
    else:
        metadata["level"] = level
    if "skill_type" in available_cols:
        columns.append("skill_type")
        values.append(skill)
    else:
        metadata["skill_type"] = skill
    if metadata_col:
        columns.append(metadata_col)
        values.append(json.dumps(metadata, ensure_ascii=False))
    columns.append("embedding")
    values.append(vector_str)

    cols_sql = ", ".join(columns)
    placeholders = ", ".join(["%s"] * len(values))
    cur.execute(f"INSERT INTO {table} ({cols_sql}) VALUES ({placeholders});", tuple(values))


def _resolve_default_sample_dir(explicit_dir: Optional[str]) -> str:
    # Ưu tiên: đối số CLI -> ENV -> mặc định theo project root
    if explicit_dir:
        return explicit_dir
    env_dir = os.getenv("RAG_SAMPLE_DATA_DIR") or os.getenv("SAMPLE_DATA_DIR")
    if env_dir:
        return env_dir
    project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    return os.path.join(project_root, "data", "sample_data")


def ingest_sample_data(
    root_dir: Optional[str] = None,
    table_name: str = "contentchunks",
) -> int:
    """
    Quét thư mục sample_data, trích xuất DOCX/PPTX, chunk, embed và chèn vào bảng vector store.
    Trả về số lượng chunk đã thêm.
    """
    root_dir = _resolve_default_sample_dir(root_dir)
    if not os.path.isdir(root_dir):
        print(f"[Ingest] Thư mục không tồn tại: {root_dir}")
        return 0

    conn = get_db_connection()
    if not conn:
        print("[Ingest] Không thể kết nối DB.")
        return 0

    inserted = 0
    try:
        with conn.cursor() as cur:
            for dirpath, _, filenames in os.walk(root_dir):
                for fname in filenames:
                    fpath = os.path.join(dirpath, fname)
                    base = os.path.basename(fpath)
                    ext = os.path.splitext(base)[1].lower()

                    level = _infer_level_from_path(fpath)
                    skill = _infer_skill_from_path(fpath)
                    meta = {
                        "doc_type": "DOCX" if ext == ".docx" else ("PPTX" if ext == ".pptx" else ext.upper()),
                        "source_path": fpath.replace("\\", "/"),
                        "tags": [],
                    }

                    if ext == ".docx":
                        lines = _read_docx(fpath)
                        chunks = _chunk_text(lines)
                        for idx, ch in enumerate(chunks, start=1):
                            _insert_chunk(cur, table_name, ch, base, idx, level, skill, meta)
                            inserted += 1
                    elif ext == ".pptx":
                        slides = _read_pptx(fpath)
                        for slide_idx, slide_text in slides:
                            # chunk theo slide nếu quá dài
                            for idx, ch in enumerate(_chunk_text([slide_text]), start=1):
                                logical_page = int(f"{slide_idx}{idx:02d}")  # ví dụ: slide 3 chunk 1 -> 301
                                _insert_chunk(cur, table_name, ch, base, logical_page, level, skill, meta)
                                inserted += 1
                    else:
                        continue

            conn.commit()
    except Exception as e:
        if conn:
            conn.rollback()
        print(f"[Ingest] Lỗi khi ingest: {e}")
    finally:
        if conn:
            conn.close()

    print(f"[Ingest] Đã thêm {inserted} chunks từ {root_dir} vào {table_name}")
    return inserted


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest sample data (DOCX/PPTX) vào bảng contentchunks")
    parser.add_argument("--dir", dest="root_dir", default=None, help="Đường dẫn thư mục sample_data. Mặc định tự suy luận theo project root hoặc ENV RAG_SAMPLE_DATA_DIR")
    parser.add_argument("--table", dest="table_name", default=os.getenv("RAG_CONTENT_CHUNK_TABLE", "contentchunks"), help="Tên bảng chunk (mặc định 'contentchunks' hoặc ENV RAG_CONTENT_CHUNK_TABLE)")
    args = parser.parse_args()
    ingest_sample_data(root_dir=args.root_dir, table_name=args.table_name)


