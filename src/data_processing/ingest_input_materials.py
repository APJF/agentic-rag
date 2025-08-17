import os
import re
import json
import argparse
from typing import List, Dict, Any, Optional, Tuple, Callable

from langchain.text_splitter import RecursiveCharacterTextSplitter
from psycopg2.extras import execute_values

from src.core.vector_store_interface import get_db_connection
from src.core.embedding import get_embedding_model
from src.config import settings


SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "txt", "md"}


# ===== File readers (tách riêng để giảm độ phức tạp) =====
def read_text_file(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def read_docx_file(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        texts: List[str] = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
        for table in getattr(doc, "tables", []) or []:
            for row in table.rows:
                row_text = " ".join([cell.text.strip() for cell in row.cells if cell.text])
                if row_text:
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception:
        return ""


def read_pptx_file(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join([run.text for run in para.runs]).strip()
                        if t:
                            texts.append(t)
                elif hasattr(shape, "text"):
                    t = str(getattr(shape, "text", "")).strip()
                    if t:
                        texts.append(t)
        return "\n".join(texts)
    except Exception:
        return ""


def read_pdf_file(file_path: str) -> str:
    try:
        import pdfplumber
        texts: List[str] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
                if t.strip():
                    texts.append(t)
        return "\n\n".join(texts)
    except Exception:
        return ""


def select_reader(ext: str) -> Optional[Callable[[str], str]]:
    ext_l = (ext or "").lower()
    if ext_l in {"txt", "md"}:
        return read_text_file
    mapping: Dict[str, Callable[[str], str]] = {
        "docx": read_docx_file,
        "pptx": read_pptx_file,
        "pdf": read_pdf_file,
    }
    return mapping.get(ext_l)


def parse_material_filename(filename: str) -> Optional[Dict[str, str]]:
    """
    Parse tên file theo quy ước:
    {COURSE}__{CHAPTER}__{UNIT}__{TYPE}__{LANG}__{SEQ}.{EXT}
    Trả None nếu không khớp.
    """
    name = os.path.basename(filename)
    match = re.match(
        r"^(?P<course>JPD\d{3})__(?P<chapter>CHAPTER_\d{2})__(?P<unit>UNIT_\d{2})__(?P<type>[A-Z_]+)__(?P<lang>[A-Z]{2}(?:_[A-Z]{2})?)__(?P<seq>\d{4})\.(?P<ext>[A-Za-z0-9]+)$",
        name,
    )
    if not match:
        return None
    parts = match.groupdict()
    parts["material_id"] = name.rsplit(".", 1)[0]
    parts["ext"] = parts["ext"].lower()
    return parts


def derive_level_from_course(course_id: str) -> Optional[str]:
    """Suy ra level từ mã course (heuristic): 1xx->N5, 2xx->N4, 3xx->N3, 4xx->N2, 5xx->N1"""
    try:
        num = int(re.sub(r"^JPD", "", course_id))
        hundreds = num // 100
        mapping = {1: "N5", 2: "N4", 3: "N3", 4: "N2", 5: "N1"}
        return mapping.get(hundreds)
    except Exception:
        return None


def extract_text_from_file(file_path: str, ext: str) -> str:
    reader = select_reader(ext)
    if not reader:
        return ""
    return reader(file_path)


def chunk_text(text: str) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100,
        length_function=len,
        add_start_index=False,
    )
    return [c.strip() for c in splitter.split_text(text) if c and len(c.strip()) >= 30]


def collect_material_files(root_dir: str) -> List[str]:
    files: List[str] = []
    for dirpath, _, filenames in os.walk(root_dir):
        for fn in filenames:
            ext = fn.rsplit(".", 1)[-1].lower() if "." in fn else ""
            if ext in SUPPORTED_EXTENSIONS:
                files.append(os.path.join(dirpath, fn))
    return files


def upsert_material_chunks(record_dicts: List[Dict[str, Any]], material_id: str, conn) -> None:
    """
    Xóa chunks cũ theo material_id (ưu tiên theo metadata_json->>'material_id'), sau đó chèn mới.
    """
    table_name = settings.RAG_CONTENT_CHUNK_TABLE

    def _get_table_columns(cursor, table: str) -> set:
        cursor.execute(
            """
            SELECT column_name FROM information_schema.columns
            WHERE table_name=%s AND table_schema='public';
            """,
            (table,)
        )
        return {r[0] for r in cursor.fetchall()}

    def _delete_existing(cursor, table: str, columns: set, mat_id: str) -> None:
        try:
            if "material_id" in columns:
                cursor.execute(f"DELETE FROM {table} WHERE material_id = %s;", (mat_id,))
            elif "metadata_json" in columns:
                cursor.execute(
                    f"DELETE FROM {table} WHERE metadata_json::jsonb ->> 'material_id' = %s;",
                    (mat_id,),
                )
        except Exception:
            pass

    def _build_insert_columns(columns: set) -> List[str]:
        cols = [
            "chunk_text",
            "embedding",
            "course_id",
            "source_document_name",
            "original_page_number",
        ]
        if "level" in columns:
            cols.append("level")
        if "skill_type" in columns:
            cols.append("skill_type")
        if "metadata_json" in columns:
            cols.append("metadata_json")
        if "material_id" in columns:
            cols.append("material_id")
        return cols

    def _build_values(records: List[Dict[str, Any]], columns: set) -> List[Tuple]:
        rows: List[Tuple] = []
        for d in records:
            base = [
                d.get("chunk_text"),
                d.get("embedding"),
                d.get("course_id"),
                d.get("source_document_name"),
                d.get("original_page_number"),
            ]
            if "level" in columns:
                base.append(d.get("level"))
            if "skill_type" in columns:
                base.append(d.get("skill_type"))
            if "metadata_json" in columns:
                base.append(d.get("metadata_json"))
            if "material_id" in columns:
                base.append(d.get("material_id"))
            rows.append(tuple(base))
        return rows

    with conn.cursor() as cur:
        cols = _get_table_columns(cur, table_name)
        _delete_existing(cur, table_name, cols, material_id)
        insert_cols = _build_insert_columns(cols)
        values = _build_values(record_dicts, cols)
        insert_query = f"INSERT INTO {table_name} ({', '.join(insert_cols)}) VALUES %s;"
        execute_values(cur, insert_query, values)
        conn.commit()


def _resolve_root_dir(args_root: str, project_root: str) -> str:
    root_dir = args_root
    if not os.path.isabs(root_dir):
        root_dir = os.path.join(project_root, root_dir)
    return os.path.normpath(root_dir)


def _build_record_dicts(
    chunks: List[str],
    fp: str,
    meta: Dict[str, str],
    course_id: str,
    chapter_id: str,
    unit_id: str,
    material_type: str,
    language: str,
    level: Optional[str],
    embedding_model,
) -> List[Dict[str, Any]]:
    material_id = meta["material_id"]
    records: List[Dict[str, Any]] = []
    for idx, chunk_text_value in enumerate(chunks, start=1):
        embedding = embedding_model.encode(chunk_text_value).tolist()
        metadata = {
            "material_id": material_id,
            "course_id": course_id,
            "chapter_id": chapter_id,
            "unit_id": unit_id,
            "material_type": material_type,
            "language": language,
            "seq": meta["seq"],
            "level": level,
            "source_path": fp.replace("\\", "/"),
            "chunk_index": idx,
        }
        records.append(
            {
                "chunk_text": chunk_text_value,
                "embedding": embedding,
                "course_id": course_id,
                "source_document_name": os.path.basename(fp),
                "original_page_number": None,
                "level": level,
                "skill_type": material_type,
                "metadata_json": json.dumps(metadata, ensure_ascii=False),
                "material_id": material_id,
            }
        )
    return records


def _process_one_file(fp: str, embedding_model, conn) -> None:
    meta = parse_material_filename(os.path.basename(fp))
    if not meta:
        print(f"BỎ QUA (không đúng quy ước tên): {fp}")
        return

    course_id = meta["course"]
    chapter_id = meta["chapter"]
    unit_id = meta["unit"]
    material_type = meta["type"]
    language = meta["lang"]
    ext = meta["ext"]

    raw_text = extract_text_from_file(fp, ext)
    if not raw_text or len(raw_text.strip()) < 30:
        print(f"BỎ QUA (không có nội dung hoặc quá ngắn): {fp}")
        return

    chunks = chunk_text(raw_text)
    if not chunks:
        print(f"BỎ QUA (không tạo được chunk): {fp}")
        return

    level = derive_level_from_course(course_id)
    records = _build_record_dicts(
        chunks=chunks,
        fp=fp,
        meta=meta,
        course_id=course_id,
        chapter_id=chapter_id,
        unit_id=unit_id,
        material_type=material_type,
        language=language,
        level=level,
        embedding_model=embedding_model,
    )
    upsert_material_chunks(records, meta["material_id"], conn)
    print(f"OK: {meta['material_id']} -> {len(records)} chunks")


def _ingest_root(root_dir: str) -> None:
    material_files = collect_material_files(root_dir)
    if not material_files:
        print(f"Không tìm thấy file đầu vào trong {root_dir}")
        return
    embedding_model = get_embedding_model()
    conn = get_db_connection()
    if not conn:
        print("Không thể kết nối DB")
        return
    try:
        print(f"Tìm thấy {len(material_files)} file trong {root_dir}")
        for fp in material_files:
            _process_one_file(fp, embedding_model, conn)
    finally:
        if conn:
            conn.close()


def main():
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    parser = argparse.ArgumentParser(description="Ingest data/input_materials vào bảng content_chunks, xóa trùng theo material_id")
    parser.add_argument("--root", default=os.path.join(project_root, "data", "input_materials"), help="Thư mục gốc input_materials")
    args = parser.parse_args()
    root_dir = _resolve_root_dir(args.root, project_root)
    if not os.path.isdir(root_dir):
        print(f"Không tìm thấy thư mục đầu vào: {root_dir}")
        return
    _ingest_root(root_dir)


if __name__ == "__main__":
    main()


